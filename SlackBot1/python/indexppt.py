import openai
from pptx import Presentation

# Chemin vers le fichier PowerPoint
chemin_fichier = "chemin_vers_votre_fichier.pptx"

# Clé d'accès à l'API OpenAI
openai.api_key = "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"

def extraire_contenu_texte(chemin_fichier):
    contenu_texte = ""
    presentation = Presentation(chemin_fichier)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        contenu_texte += run.text
    return contenu_texte

def generer_reponse(requete_utilisateur, base_de_connaissances):
    prompt = f"{base_de_connaissances}\nQuestion : {requete_utilisateur}\nRéponse :"
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=100,
        temperature=0.7,
        n=1,
        stop=None
    )
    reponse = response.choices[0].text.strip()
    return reponse

def poser_question():
    # Capturer la requête de l'utilisateur
    requete_utilisateur = input("Posez votre question : ")

    # Extraire le contenu texte du fichier PowerPoint
    contenu_texte = extraire_contenu_texte(chemin_fichier)

    # Générer et afficher la réponse
    reponse = generer_reponse(requete_utilisateur, contenu_texte)
    print("Réponse : ", reponse)

def main():
    while True:
        poser_question()

if __name__ == '__main__':
    main()
